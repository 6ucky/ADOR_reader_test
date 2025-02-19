from llama_index.core import SimpleDirectoryReader
from llama_index.readers.pdf_table import PDFTableReader
import docx2txt
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import PyPDF2
import textwrap
from PIL import Image as PILImage
from io import BytesIO
from pathlib import Path

from llama_index.core.schema import Document

import os
import fitz
import io

class ChunkedDocument(Document):
    @classmethod
    def from_document(cls, doc: Document, chunk_size: int = 1000, overlap: int = 100):
        chunks = textwrap.wrap(doc.text, chunk_size, break_long_words=False, replace_whitespace=False)
        chunked_docs = []
        
        for i, chunk in enumerate(chunks):
            if i > 0:
                # Add overlap from previous chunk
                chunk = chunks[i-1][-overlap:] + chunk
            if i < len(chunks) - 1:
                # Add overlap to next chunk
                chunk = chunk + chunks[i+1][:overlap]
            
            chunked_doc = cls(text=chunk, extra_info=doc.extra_info.copy())
            chunked_doc.extra_info['chunk_id'] = i
            chunked_docs.append(chunked_doc)
        
        return chunked_docs

class DocxReader:
    def load_data(self, file, extra_info=None):
        try:
            with open(file, 'rb') as f:
                doc = DocxDocument(f)
                full_text = '\n'.join(paragraph.text for paragraph in doc.paragraphs if paragraph.text)
                metadata = {
                    'title': doc.core_properties.title,
                    'filename': os.path.basename(file),
                    'author': doc.core_properties.author,
                    'keywords': doc.core_properties.keywords,
                    'created': doc.core_properties.created.isoformat() if doc.core_properties.created else None,
                    'modified': doc.core_properties.modified.isoformat() if doc.core_properties.modified else None,
                    'type': 'text'
                }
                all_documents = [Document(text=full_text, extra_info=metadata)]

                # extract tables from docs 
                metadata['type'] = 'table'
                for table in doc.tables:
                    table_content = []
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_content.append(','.join(row_data))
                    all_documents.append(Document(text='\n'.join(table_content), extra_info=metadata))

                # extract images from docx
                rels = doc.part.rels
                for rel in rels:
                    if "image" in rels[rel].target_ref:
                        img = rels[rel].target_part
                        img_data = img.blob
                        file_name = os.path.basename(file)
                        img_filename = file_name + '_' + os.path.basename(img.partname)
                        image = Image.open(BytesIO(img_data))
                        if image.width >= 100 and image.height >= 100:
                            with open(os.path.join(image_folder, img_filename), "wb") as f:
                                f.write(img_data)

                return all_documents
        except Exception as e:
            print(f"Erreur lors de la lecture du fichier DOCX {file}: {e}")
            pass

class XlsxReader:
    def load_data(self, file, extra_info=None):
        try:
            with open(file, 'rb') as f:
                wb = openpyxl.load_workbook(f)
                text = []
                for sheet in wb:
                    for row in sheet.iter_rows(values_only=True):
                        text.append(' '.join(str(cell) for cell in row if cell is not None))
                metadata = {
                    'title': wb.properties.title,
                    'filename': os.path.basename(file),
                    'creator': wb.properties.creator,
                    'created': wb.properties.created.isoformat() if wb.properties.created else None,
                    'modified': wb.properties.modified.isoformat() if wb.properties.modified else None,
                    'sheet_names': wb.sheetnames,
                    'type': 'text'
                }
                # extract images
                # TODO: wmf image can not be extracted
                image_count = 0
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    for image in sheet._images:
                        img_bytes = image._data()
                        img_pil = PILImage.open(io.BytesIO(img_bytes))
                        image_path = os.path.join(image_folder, f'{os.path.basename(file)}_image{image_count}.png')
                        img_pil.save(image_path)
                        image_count += 1

                return [Document(text='\n'.join(text), extra_info=metadata)]
        except Exception as e:
            print(f"Erreur lors de la lecture du fichier XLSX {file}: {e}")
            pass

class PptxReader:
    def load_data(self, file, extra_info=None):
        try:
            with open(file, 'rb') as f:
                prs = Presentation(f)
                full_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            full_text.append(shape.text)
                metadata = {
                    'title': prs.core_properties.title,
                    'filename': os.path.basename(file),
                    'author': prs.core_properties.author,
                    'created': prs.core_properties.created.isoformat() if prs.core_properties.created else None,
                    'modified': prs.core_properties.modified.isoformat() if prs.core_properties.modified else None,
                    'slides': len(prs.slides),
                    'type': 'text'
                }
                all_documents = [Document(text='\n'.join(full_text), extra_info=metadata)]

                image_count = 0
                for slide in prs.slides:
                    text_forms = ''
                    for shape in slide.shapes:
                        # extract smartArt
                        if shape.shape_type == 1:
                            text = shape.text_frame.text
                            metadata['type'] = 'smartart'
                            all_documents.append(Document(text=text, extra_info=metadata))
                        # extract tables
                        elif shape.has_table:
                            table = shape.table
                            table_data = []
                            for row in table.rows:
                                row_data = [cell.text for cell in row.cells]
                                table_data.append(','.join(row_data))
                            metadata['type'] = 'table'
                            all_documents.append(Document(text='\n'.join(table_data), extra_info=metadata))
                        # extract text from forms
                        elif shape.has_text_frame:
                            text_forms += shape.text_frame.text + '\n'
                        # extract images
                        elif shape.shape_type == 13:
                            image = shape.image
                            image_bytes = image.blob
                            image_filename = f'{os.path.basename(file)}_image{image_count}.{image.ext}'
                            image_path = os.path.join(image_folder, image_filename)
                            with open(image_path, 'wb') as f:
                                f.write(image_bytes)
                            image_count += 1
                    if len(text_forms) > 30:
                        metadata['type'] = 'forms'
                        all_documents.append(Document(text=text_forms, extra_info=metadata))

                return all_documents
        except Exception as e:
            print(f"Erreur lors de la lecture du fichier PPTX {file}: {e}")
            pass


class PdfReader:
    def load_data(self, file, extra_info=None):
        try:
            with open(file, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
                metadata = {
                    'filename': os.path.basename(file),
                    'num_pages': len(pdf_reader.pages),
                    'type': 'text'
                }

                # camelot to extract tables
                metadata['type'] = 'table'
                reader = PDFTableReader()
                pdf_path = Path(file)
                documents = reader.load_data(file=pdf_path, extra_info=metadata)
                documents.append(Document(text=text, extra_info=metadata))

                # extract images
                self._extract_images(file)

                return documents
        except Exception as e:
            print(f"Erreur lors de la lecture du fichier PDF {file}: {e}")
            pass

    def _extract_images(self, file):
        pdf = fitz.open(file)
        file_name = os.path.basename(file)
        page_number = 0
        for page in pdf:
            page_number += 1
            image_number = 0
            for block in page.get_text("dict")["blocks"]:
                # Skip if not an image block
                if block["type"] != 1:
                    continue
                # skip if image mostly black
                if self.is_mostly_black(block["image"]) == False:
                    continue
                # skip if a span window or a small icon
                if block['width'] >= 100 and block['height'] >= 100:
                    image_number += 1
                    pix = page.get_pixmap(dpi=300, clip=block['bbox'])
                    output_file_name = f"{image_folder}/{file_name}_image{image_number}.png"
                    pix.pil_save(output_file_name)

    def is_mostly_black(self, image_data, threshold=0.8):
        # Load the image
        img = Image.open(io.BytesIO(image_data))

        # Convert to grayscale
        img_gray = img.convert("L")

        # Calculate the number of non-black pixels
        num_non_black = sum(1 for pixel in img_gray.getdata() if pixel > 0)

        # Calculate the total number of pixels
        total_pixels = img_gray.width * img_gray.height

        # Calculate the percentage of non-black pixels
        percentage_non_black = num_non_black / total_pixels

        # Check if the percentage exceeds the threshold
        return percentage_non_black >= threshold
